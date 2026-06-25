"""Ingestion service — the document → vectors pipeline (CLAUDE.md §4, §5).

This is the research core of the ingestion side. It is deliberately split into
small, single-responsibility functions:

    parse_*          bytes → list of {text, page} dicts   (one per parser)
    parse_document   routes to the right parser by extension
    hierarchical_chunk   pages → child/parent chunks      (PURE — no I/O)
    embed_chunks     texts → 768-d vectors                (Gemini)
    upsert_to_qdrant chunks + vectors → Qdrant point ids
    process_document orchestrates all of the above        (background task)

The flow:  parse → clean → hierarchical-chunk → embed children → upsert to
Qdrant (+payload) → persist chunk rows + mark the document ready.

`hierarchical_chunk` is intentionally pure (no DB / Gemini / Qdrant calls) so the
chunking logic — the part most worth getting right and testing — can be exercised
in isolation. All model names / dims / batch sizes come from settings, never
hardcoded.
"""

import asyncio
import hashlib
import io
import logging
import re
import uuid

import fitz  # PyMuPDF — keeps page numbers, which our citations depend on.
import google.generativeai as genai
from docx import Document as DocxDocument
from pptx import Presentation
from qdrant_client.models import FieldCondition, Filter, FilterSelector, MatchValue, PointStruct, SparseVector
from sqlalchemy import select

from config import get_settings
from database import get_session
from models.document import Document, DocumentChunk
from qdrant import get_qdrant_client

logger = logging.getLogger(__name__)

# ── Chunk sizing (CLAUDE.md §5) ───────────────────────────────────────────────
# We approximate tokens with words because exact tokenization is model-specific
# and not worth the dependency here: 1 token ≈ 0.75 words, so N tokens ≈ N/0.75
# words. Child = 256 tokens ≈ 341 words; Parent = 1024 tokens ≈ 1365 words.
CHILD_WORDS = 341          # ~256 tokens — embedded + retrieved (precision)
PARENT_WORDS = 1365        # ~1024 tokens — sent to the LLM (context)
CHILD_OVERLAP_WORDS = 68   # 20% of CHILD_WORDS — overlap so an answer that
                           # straddles a boundary survives inside one child.

# Task type for the embedding API. This is the DOCUMENT side of an asymmetric
# embedding model — the query path (a later session) must use "retrieval_query"
# instead, so this constant deliberately lives here and is NOT shared.
EMBED_TASK_DOCUMENT = "retrieval_document"


# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────
def _clean_text(text: str) -> str:
    # Strip *excessive* whitespace while preserving paragraph breaks (a blank
    # line). Runs of spaces/tabs collapse to one space; 3+ newlines collapse to a
    # double newline (one paragraph break); trailing spaces per line are dropped.
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ──────────────────────────────────────────────────────────────────────────────
# Boilerplate stripping — remove repeated headers/footers/letterhead.
#
# WHY: a fixed letterhead, classification banner, and "Page X of Y" footer repeat
# on every page. Because every chunk then starts with the SAME text, chunks from
# different documents look semantically similar — diluting embeddings and blunting
# the reranker's ability to tell relevant context apart. Stripping this boilerplate
# improves both retrieval precision and reranker discrimination.
# ──────────────────────────────────────────────────────────────────────────────

# Company name on the letterhead (the corpus is "INFOVANCE TECHNOLOGIES").
_COMPANY_RE = re.compile(r"infovance", re.IGNORECASE)
# Document classification / distribution banners.
_CLASSIFICATION_RE = re.compile(
    r"^(classification\s*[:\-]|confidential|internal use only|internal\b|"
    r"restricted|public|do not distribute)",
    re.IGNORECASE,
)
# Bare page markers: "Page 3", "Page 3 of 12", "3", "- 4 -".
_PAGENUM_RE = re.compile(
    r"^(page\s+)?\d+(\s+of\s+\d+)?$|^[-–—]\s*\d+\s*[-–—]$", re.IGNORECASE
)
# Footer lines ending in a page marker, e.g. "HR-POL-LV-2025-06 - v3.2 - Page 1".
# The trailing page number varies per page, so the frequency pass can't catch these.
_FOOTER_RE = re.compile(r"\bpage\s+\d+(\s+of\s+\d+)?\s*$", re.IGNORECASE)


def _looks_like_label(line: str) -> bool:
    # A short (<4-word) line that reads as a structural label/header, not prose:
    # all-caps (e.g. a doc code "HR-POL-LV-2025-06" or "CASUAL LEAVE") or a
    # colon-terminated key ("Owner:"). Sentence-case content ("12 days") is kept.
    if len(line.split()) >= 4:
        return False
    letters = [c for c in line if c.isalpha()]
    all_caps = bool(letters) and all(c.isupper() for c in letters)
    return all_caps or line.rstrip().endswith(":")


def _is_pattern_boilerplate(line: str) -> bool:
    # Per-line boilerplate test, applied regardless of how often the line repeats
    # (footers with page numbers vary per page, so frequency alone won't catch them).
    s = line.strip()
    if not s:
        return True
    return bool(
        # .match (anchored at start): strips letterhead/footer lines that BEGIN with
        # the company name, but leaves prose that merely mentions it mid-sentence.
        _COMPANY_RE.match(s)
        or _CLASSIFICATION_RE.match(s)
        or _PAGENUM_RE.match(s)
        or _FOOTER_RE.search(s)
        or _looks_like_label(s)
    )


def _strip_boilerplate(pages: list[dict]) -> list[dict]:
    """Remove repeated header/footer/letterhead lines from already-parsed pages.

    Two passes: (1) drop lines that appear on MORE than 60% of pages (the recurring
    letterhead/footer), and (2) drop lines matching boilerplate patterns (company
    name, classification banner, bare page number, short all-caps/colon label)
    regardless of frequency. Pages that become empty are dropped.
    """
    n = len(pages)
    # Pass 1: count how many pages each (stripped) line appears on. `set(...)` per
    # page so a line repeated within one page is counted once.
    freq: dict[str, int] = {}
    for p in pages:
        for line in {ln.strip() for ln in p["text"].split("\n") if ln.strip()}:
            freq[line] = freq.get(line, 0) + 1
    # >60% of pages = structural boilerplate. Needs ≥2 pages to tell repetition from
    # a single-page document (where every line trivially appears on "100%" of pages).
    repeated = {line for line, c in freq.items() if c / n > 0.6} if n >= 2 else set()

    cleaned: list[dict] = []
    for p in pages:
        kept = [
            ln
            for ln in p["text"].split("\n")
            if ln.strip()
            and ln.strip() not in repeated
            and not _is_pattern_boilerplate(ln)
        ]
        text = "\n".join(kept).strip()
        if text:
            cleaned.append({"text": text, "page": p["page"]})
    return cleaned


# ──────────────────────────────────────────────────────────────────────────────
# Function 1–3: format-specific parsers (bytes → list[{text, page}])
# ──────────────────────────────────────────────────────────────────────────────
def parse_pdf(file_bytes: bytes) -> list[dict]:
    """Extract text per page from a PDF using PyMuPDF (fitz)."""
    pages: list[dict] = []
    # Open straight from bytes — the upload never touches disk.
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    try:
        for i, page in enumerate(doc):
            text = _clean_text(page.get_text())
            # Skip pages with no extractable text — these are scanned/image-only
            # pages. We silently drop them; OCR here would be future work.
            if not text:
                continue
            # 1-based page number: it's what a human reads on the page and what a
            # citation ("page 3") must show. We keep it because citations need it.
            pages.append({"text": text, "page": i + 1})
    finally:
        doc.close()
    # Strip the recurring letterhead/footer now that we have every page to compare.
    return _strip_boilerplate(pages)


def parse_docx(file_bytes: bytes) -> list[dict]:
    """Extract every paragraph from a .docx using python-docx."""
    document = DocxDocument(io.BytesIO(file_bytes))
    pages: list[dict] = []
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue  # skip empty paragraphs
        # Drop short header/metadata paragraphs (doc codes, classification labels,
        # all-caps section labels) — same boilerplate test the PDF path uses.
        if _is_pattern_boilerplate(text):
            continue
        # Word has no reliable page-number API (pagination is computed by the
        # renderer, not stored), so page is None — citations from a DOCX show the
        # filename only, no page.
        pages.append({"text": text, "page": None})
    return pages


def parse_pptx(file_bytes: bytes) -> list[dict]:
    """Extract title + body text per slide from a .pptx using python-pptx."""
    prs = Presentation(io.BytesIO(file_bytes))
    pages: list[dict] = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []
        # Every text-bearing shape (the title placeholder is one of these) — so
        # this captures both the slide title and all body text frames.
        for shape in slide.shapes:
            if shape.has_text_frame:
                frame_text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                if frame_text.strip():
                    parts.append(frame_text)
        text = _clean_text("\n".join(parts))
        if not text:
            continue
        # Slide number (1-based) is used as the "page" for citations.
        pages.append({"text": text, "page": i + 1})
    # Strip recurring slide headers/footers + per-line label boilerplate.
    return _strip_boilerplate(pages)


# ──────────────────────────────────────────────────────────────────────────────
# Function 4: single entry point — routes by extension
# ──────────────────────────────────────────────────────────────────────────────
def parse_document(file_bytes: bytes, filename: str) -> list[dict]:
    """Route to the correct parser by file extension.

    Single entry point so callers never branch on file type themselves.
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pdf":
        return parse_pdf(file_bytes)
    if ext == "docx":
        return parse_docx(file_bytes)
    if ext == "pptx":
        return parse_pptx(file_bytes)
    raise ValueError(f"Unsupported file type: .{ext}. Allowed: pdf, docx, pptx")


# ──────────────────────────────────────────────────────────────────────────────
# Function 5: hierarchical chunking — the most important function (PURE)
# ──────────────────────────────────────────────────────────────────────────────
def hierarchical_chunk(pages: list[dict]) -> list[dict]:
    """Turn parsed pages into hierarchical child/parent chunks.

    WHY TWO SIZES (CLAUDE.md §5):
      • Child (~256 tokens): small, so a vector match is precise — the retrieved
        text is tightly about the query. This is what we embed and search on.
      • Parent (~1024 tokens): large, so when we hand text to the LLM it has
        enough surrounding context to actually answer well. We retrieve on the
        child but generate from its parent.
      • Child overlap (20%): a fact split across a child boundary would otherwise
        be cut in half; the overlap guarantees it stays whole in at least one
        child. Parents do NOT overlap (they're context blocks, not match units).

    Structure: the document is split into parents first, then each parent is
    sub-split into overlapping children that carry that parent's full text. So a
    retrieved child always maps back to exactly one parent.

    PAGE TRACKING: we flatten all pages into one word stream while remembering
    each word's source page, so every child reports the page it STARTS on. That
    is what makes "[Source: file.pdf, page 3]" citations accurate.

    Pure function: no DB / Gemini / Qdrant. Easy to unit-test on its own.
    """
    # Flatten every page to a flat word list + a parallel page list (word i came
    # from page word_pages[i]). word-level granularity is what lets a child know
    # its starting page after the page boundaries are gone.
    words: list[str] = []
    word_pages: list[int | None] = []
    for page in pages:
        for w in page["text"].split():
            words.append(w)
            word_pages.append(page["page"])

    # Nothing extractable (e.g. an all-image scanned PDF): no chunks.
    if not words:
        return []

    # Edge case: the whole document is shorter than a single child chunk. Don't
    # emit an empty/duplicate parent — use the full text as BOTH child and parent.
    if len(words) <= CHILD_WORDS:
        full_text = " ".join(words)
        return [
            {
                "child_text": full_text,
                "parent_text": full_text,
                "page_number": word_pages[0],
                "chunk_index": 0,
            }
        ]

    chunks: list[dict] = []
    chunk_index = 0
    child_step = CHILD_WORDS - CHILD_OVERLAP_WORDS  # how far to advance each child

    # Outer loop: parents (1024 tokens, NO overlap).
    parent_start = 0
    while parent_start < len(words):
        parent_end = min(parent_start + PARENT_WORDS, len(words))
        parent_word_slice = words[parent_start:parent_end]
        parent_page_slice = word_pages[parent_start:parent_end]
        parent_text = " ".join(parent_word_slice)

        # Inner loop: children (256 tokens, 20% overlap) carved out of THIS parent.
        child_start = 0
        while child_start < len(parent_word_slice):
            child_end = min(child_start + CHILD_WORDS, len(parent_word_slice))
            child_text = " ".join(parent_word_slice[child_start:child_end])
            chunks.append(
                {
                    "child_text": child_text,
                    "parent_text": parent_text,
                    # Page where this child starts; None for DOCX.
                    "page_number": parent_page_slice[child_start],
                    "chunk_index": chunk_index,
                }
            )
            chunk_index += 1
            # Stop once a child reaches the parent's end — otherwise the overlap
            # step would emit a tiny trailing duplicate (or loop forever).
            if child_end == len(parent_word_slice):
                break
            child_start += child_step

        parent_start += PARENT_WORDS  # next parent — no overlap between parents

    return chunks


# ──────────────────────────────────────────────────────────────────────────────
# Function 6: embed child texts with Gemini (768-d, batched)
# ──────────────────────────────────────────────────────────────────────────────
async def embed_chunks(texts: list[str]) -> list[list[float]]:
    """Embed texts with Gemini, returning one 768-d vector per input.

    • model + dims come from settings (gemini-embedding-001 @ 768d).
    • output_dimensionality=768 is Matryoshka truncation of the model's native
      3072-d output — same quality, but smaller Qdrant storage and faster
      similarity search. (No re-normalization needed: the collection uses cosine
      distance, which is scale-invariant.)
    • Batched in groups of EMBEDDING_BATCH_SIZE with a 1s pause between batches to
      stay under Gemini's free-tier RPM limit.
    """
    if not texts:
        return []

    settings = get_settings()
    genai.configure(api_key=settings.GEMINI_API_KEY)
    # The SDK expects a fully-qualified "models/<name>" id; settings stores the
    # bare name, so prefix it if the caller hasn't.
    model_name = settings.GEMINI_EMBEDDING_MODEL
    if not model_name.startswith("models/"):
        model_name = f"models/{model_name}"

    batch_size = settings.EMBEDDING_BATCH_SIZE
    embeddings: list[list[float]] = []

    for start in range(0, len(texts), batch_size):
        batch = texts[start : start + batch_size]
        # genai.embed_content is a blocking (sync) HTTP call — run it in a thread
        # so it doesn't stall the event loop while the background task waits.
        result = await asyncio.to_thread(
            genai.embed_content,
            model=model_name,
            content=batch,
            task_type=EMBED_TASK_DOCUMENT,
            output_dimensionality=settings.GEMINI_EMBEDDING_DIMENSIONS,
        )
        # With a list `content`, the SDK returns {"embedding": [[...], [...], ...]}.
        embeddings.extend(result["embedding"])

        # Pause between batches (but not after the final one) to respect RPM.
        # asyncio.sleep (not time.sleep) so we yield the loop instead of blocking.
        if start + batch_size < len(texts):
            await asyncio.sleep(1)

    return embeddings


# ──────────────────────────────────────────────────────────────────────────────
# Function 7: sparse vector builder + upsert helpers
# ──────────────────────────────────────────────────────────────────────────────
def build_sparse_vector(text: str) -> SparseVector:
    """Build a BM25-approximate sparse vector (TF-weighted, hash-indexed).

    Index = MD5(word) % 100_000; value = term frequency. We use hashlib.md5
    instead of Python's built-in hash() because hash() is randomized per process
    (PYTHONHASHSEED). A random seed means the same word gets a different index on
    every restart, so query and document sparse vectors never overlap. MD5 is
    stable, fast, and well-distributed — correct for a research BM25 approximation.
    Production would use SPLADE or FastEmbed BM25.
    """
    tokens = re.sub(r"[^a-z0-9\s]", " ", text.lower()).split()
    if not tokens:
        return SparseVector(indices=[], values=[])
    total = len(tokens)
    counts: dict[str, int] = {}
    for tok in tokens:
        counts[tok] = counts.get(tok, 0) + 1
    index_values: dict[int, float] = {}
    for word, count in counts.items():
        idx = int(hashlib.md5(word.encode()).hexdigest(), 16) % 100_000
        index_values[idx] = index_values.get(idx, 0.0) + count / total
    indices = list(index_values.keys())
    values = [index_values[i] for i in indices]
    return SparseVector(indices=indices, values=values)


def _batch_upsert_points(points: list[PointStruct]) -> None:
    """Sync helper: upsert Qdrant points in batches of 50."""
    settings = get_settings()
    client = get_qdrant_client()
    for start in range(0, len(points), 50):
        client.upsert(
            collection_name=settings.QDRANT_COLLECTION,
            points=points[start : start + 50],
        )


def upsert_to_qdrant(
    chunks: list[dict],
    embeddings: list[list[float]],
    document_id: uuid.UUID,
    department_id: uuid.UUID,
    filename: str,
    visibility: str,
    uploaded_by: uuid.UUID,
) -> list[str]:
    """Upsert one Qdrant point per chunk (dense + sparse); return the point ids.

    visibility + uploaded_by are written into EVERY point's payload so the
    3-tier visibility filter (build_visibility_filter) can enforce access at
    query time without any DB round-trip (CLAUDE.md §6).
    """
    points: list[PointStruct] = []
    point_ids: list[str] = []
    for chunk, embedding in zip(chunks, embeddings):
        point_id = str(uuid.uuid4())
        point_ids.append(point_id)
        points.append(
            PointStruct(
                id=point_id,
                vector={
                    "dense": embedding,
                    "sparse": build_sparse_vector(chunk["child_text"]),
                },
                payload={
                    "document_id": str(document_id),
                    "department_id": str(department_id),
                    "visibility": visibility,
                    "uploaded_by": str(uploaded_by),
                    "chunk_index": chunk["chunk_index"],
                    "parent_text": chunk["parent_text"],
                    "source_filename": filename,
                    "page_number": chunk["page_number"],
                },
            )
        )
    _batch_upsert_points(points)
    return point_ids


def delete_document_vectors(document_id: uuid.UUID) -> None:
    """Delete every Qdrant point whose payload.document_id matches this doc.

    Sync + filter-based: idempotent (safe to call even if vectors were never
    upserted). Run via asyncio.to_thread so the event loop stays free.
    """
    settings = get_settings()
    client = get_qdrant_client()
    client.delete(
        collection_name=settings.QDRANT_COLLECTION,
        points_selector=FilterSelector(
            filter=Filter(must=[
                FieldCondition(key="document_id", match=MatchValue(value=str(document_id)))
            ])
        ),
    )


async def reindex_document(
    document_id: uuid.UUID,
    department_id: uuid.UUID,
    filename: str,
    visibility: str,
    uploaded_by: uuid.UUID,
) -> int:
    """Re-embed all chunks and upsert with dense+sparse vectors.

    Uses existing qdrant_point_ids from DocumentChunk rows so the upsert
    updates in-place rather than creating duplicate points. Re-writes the
    visibility + uploaded_by payload so reindex also backfills these fields on
    documents ingested before the 3-tier model existed.
    Returns the number of chunks re-indexed.
    """
    async with get_session() as db:
        result = await db.execute(
            select(DocumentChunk)
            .where(DocumentChunk.document_id == document_id)
            .order_by(DocumentChunk.chunk_index)
        )
        chunks = result.scalars().all()

    if not chunks:
        return 0

    child_texts = [c.child_text for c in chunks]
    embeddings = await embed_chunks(child_texts)

    points: list[PointStruct] = []
    for chunk, embedding in zip(chunks, embeddings):
        point_id = chunk.qdrant_point_id or str(uuid.uuid4())
        points.append(
            PointStruct(
                id=point_id,
                vector={
                    "dense": embedding,
                    "sparse": build_sparse_vector(chunk.child_text),
                },
                payload={
                    "document_id": str(document_id),
                    "department_id": str(department_id),
                    "visibility": visibility,
                    "uploaded_by": str(uploaded_by),
                    "chunk_index": chunk.chunk_index,
                    "parent_text": chunk.parent_text,
                    "source_filename": filename,
                    "page_number": chunk.page_number,
                },
            )
        )

    await asyncio.to_thread(_batch_upsert_points, points)
    return len(points)


# ──────────────────────────────────────────────────────────────────────────────
# Function 8: orchestrator — runs as a FastAPI BackgroundTask
# ──────────────────────────────────────────────────────────────────────────────
async def process_document(
    file_bytes: bytes,
    filename: str,
    original_filename: str,
    document_id: uuid.UUID,
    department_id: uuid.UUID,
    visibility: str,
    uploaded_by: uuid.UUID,
) -> None:
    """Run the full pipeline for one uploaded document, in the background.

    Runs OUTSIDE a request context, so it owns its own DB session via
    get_session() (NOT get_db, whose session is tied to the finished request).
    The router commits the Document row before scheduling this, so the row is
    already visible when we load it here.

    On ANY failure we mark the document failed and swallow the exception — a
    background task must never crash the server.
    """
    logger.info("Ingestion start: doc=%s file=%s", document_id, original_filename)
    try:
        # 1. Parse → pages.
        pages = parse_document(file_bytes, filename)
        logger.info("Parsed %s → %d page(s) with text", original_filename, len(pages))

        # 2. Chunk → child/parent chunks (pure, no I/O).
        chunks = hierarchical_chunk(pages)
        if not chunks:
            # No extractable text (e.g. a scanned, image-only PDF). Fail loudly
            # with a clear message rather than silently marking an empty doc ready.
            raise ValueError(
                "No extractable text found — the file may be scanned/image-only."
            )
        logger.info("Chunked %s → %d child chunk(s)", original_filename, len(chunks))

        # 3. Embed the CHILD text only (children are the retrieval unit).
        child_texts = [c["child_text"] for c in chunks]
        embeddings = await embed_chunks(child_texts)
        logger.info("Embedded %d chunk(s) @ %dd", len(embeddings), get_settings().GEMINI_EMBEDDING_DIMENSIONS)

        # 4. Upsert vectors + payload to Qdrant. The client is sync, so run it in
        #    a thread to keep the event loop free.
        point_ids = await asyncio.to_thread(
            upsert_to_qdrant,
            chunks,
            embeddings,
            document_id,
            department_id,
            filename,
            visibility,
            uploaded_by,
        )
        logger.info("Upserted %d point(s) to Qdrant", len(point_ids))

        # Highest page number we saw text on — a good proxy for page_count, and
        # None for DOCX (no page numbers at all).
        page_numbers = [p["page"] for p in pages if p["page"] is not None]
        page_count = max(page_numbers) if page_numbers else None

        # 5 + 6. Persist chunk rows and flip the document to "ready", in one
        #        transaction owned by THIS task's session.
        async with get_session() as session:
            for chunk, point_id in zip(chunks, point_ids):
                session.add(
                    DocumentChunk(
                        document_id=document_id,
                        chunk_index=chunk["chunk_index"],
                        child_text=chunk["child_text"],
                        parent_text=chunk["parent_text"],
                        page_number=chunk["page_number"],
                        qdrant_point_id=point_id,
                    )
                )
            doc = await session.get(Document, document_id)
            if doc is not None:
                doc.status = "ready"
                doc.chunk_count = len(chunks)
                doc.page_count = page_count
                # updated_at refreshes automatically via the model's onupdate.

        logger.info(
            "Ingestion done: doc=%s status=ready chunks=%d pages=%s",
            document_id, len(chunks), page_count,
        )

    except Exception as exc:
        # Log the FULL traceback for debugging...
        logger.exception("Ingestion failed: doc=%s file=%s", document_id, original_filename)
        # ...then record the failure on the document in a SEPARATE session (the
        # one above may already be rolled back). Truncate to fit the column.
        try:
            async with get_session() as session:
                doc = await session.get(Document, document_id)
                if doc is not None:
                    doc.status = "failed"
                    doc.error_message = str(exc)[:1000]
        except Exception:
            logger.exception("Could not mark doc=%s failed", document_id)
        # Do NOT re-raise — a background task must not take the server down.
